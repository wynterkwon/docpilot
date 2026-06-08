from __future__ import annotations

from pathlib import Path

from docpilot.exceptions import IngestionError
from docpilot.ingestion.models import IngestedDocument


def ingest(path: str | Path) -> IngestedDocument:
    path = Path(path)

    if not path.exists():
        raise IngestionError("File not found", detail=str(path))

    if path.suffix.lower() != ".pptx":
        raise IngestionError(f"Expected .pptx, got '{path.suffix}'", detail=str(path))

    try:
        from pptx import Presentation
    except ImportError as e:
        raise IngestionError("python-pptx is required: pip install python-pptx") from e

    try:
        prs = Presentation(str(path))
    except Exception as e:
        raise IngestionError("Failed to open PPTX", detail=str(e)) from e

    slides = [_extract_slide(i + 1, slide) for i, slide in enumerate(prs.slides)]
    content = "\n\n".join(slides)

    return IngestedDocument(
        source=path,
        content=content,
        mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        metadata={
            "slide_count": len(prs.slides),
            "size_bytes": path.stat().st_size,
        },
    )


def _extract_slide(index: int, slide) -> str:
    parts: list[str] = [f"[슬라이드 {index}]"]

    title = _get_title(slide)
    if title:
        parts.append(f"제목: {title}")

    body_lines = _get_body(slide, title)
    if body_lines:
        parts.append("\n".join(body_lines))

    notes = _get_notes(slide)
    if notes:
        parts.append(f"노트: {notes}")

    return "\n".join(parts)


def _get_title(slide) -> str:
    if slide.shapes.title and slide.shapes.title.has_text_frame:
        return slide.shapes.title.text.strip()
    return ""


def _get_body(slide, title: str) -> list[str]:
    lines: list[str] = []
    for shape in slide.shapes:
        if shape == slide.shapes.title:
            continue
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if not text or text == title:
                continue
            level = para.level  # 0 = 최상위, 1+ = 들여쓰기
            indent = "  " * level
            prefix = "- " if level > 0 else ""
            lines.append(indent + prefix + text)
    return lines


def _get_notes(slide) -> str:
    if slide.has_notes_slide:
        notes_tf = slide.notes_slide.notes_text_frame
        return notes_tf.text.strip()
    return ""
