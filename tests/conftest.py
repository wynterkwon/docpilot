from __future__ import annotations

import io
import zipfile
from pathlib import Path

import pytest

# Minimal HWPX content.hml XML with one heading and one placeholder
_HWPX_CONTENT = """<?xml version="1.0" encoding="UTF-8"?>
<hml xmlns:hp="http://www.hancom.co.kr/hwpml/2012/paragraph"
     xmlns:hc="http://www.hancom.co.kr/hwpml/2012/core">
  <hp:body>
    <hp:sec>
      <hp:p>
        <hp:pPr>
          <hp:rPr>
            <hp:sz hp:val="2000"/>
            <hp:b hp:val="true"/>
          </hp:rPr>
        </hp:pPr>
        <hp:t>서론</hp:t>
      </hp:p>
      <hp:p>
        <hp:t>{{서론}}</hp:t>
      </hp:p>
      <hp:p>
        <hp:pPr>
          <hp:rPr>
            <hp:sz hp:val="1800"/>
            <hp:b hp:val="true"/>
          </hp:rPr>
        </hp:pPr>
        <hp:t>결론</hp:t>
      </hp:p>
      <hp:p>
        <hp:t>{{결론}}</hp:t>
      </hp:p>
    </hp:sec>
  </hp:body>
</hml>""".encode("utf-8")


def _make_hwpx(path: Path, content_xml: bytes = _HWPX_CONTENT) -> Path:
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("mimetype", "application/hwp+zip", compress_type=zipfile.ZIP_STORED)
        zf.writestr("Contents/content.hml", content_xml)
        zf.writestr("META-INF/container.xml", b"<container/>")
    return path


@pytest.fixture()
def sample_hwpx(tmp_path: Path) -> Path:
    return _make_hwpx(tmp_path / "sample.hwpx")


@pytest.fixture()
def sample_txt(tmp_path: Path) -> Path:
    p = tmp_path / "sample.txt"
    p.write_text("2025년 사업 계획서\n\n핵심 목표: 매출 30% 성장\n세부 계획: ...", encoding="utf-8")
    return p


@pytest.fixture()
def sample_pptx(tmp_path: Path) -> Path:
    pptx = pytest.importorskip("pptx")
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "2025 사업 계획"
    slide.placeholders[1].text = "핵심 내용"

    path = tmp_path / "sample.pptx"
    prs.save(str(path))
    return path


@pytest.fixture()
def hwpx_template(tmp_path: Path) -> Path:
    return _make_hwpx(tmp_path / "template.hwpx")


@pytest.fixture()
def make_hwpx(tmp_path: Path):
    """Factory fixture: make_hwpx(name, xml=...) -> Path"""
    def _factory(name: str = "doc.hwpx", content_xml: bytes = _HWPX_CONTENT) -> Path:
        return _make_hwpx(tmp_path / name, content_xml)
    return _factory
